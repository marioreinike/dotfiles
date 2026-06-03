"""
Helpers para construir presentaciones desde el template oficial de Vambe 2026.

Flujo recomendado:
    from pptx import Presentation
    import pptx_helpers as vb

    prs = Presentation("assets/Plantilla-PPT-Vambe-2026.pptx")
    vb.describe(prs)                      # ver índices de slides y shapes (introspección)
    vb.keep_only(prs, [0, 2, 5, 34])      # quedarse solo con portada, sección, métricas, cierre
    # ...o clonar un layout para repetirlo:
    nuevo = vb.clone_slide(prs, 3)        # duplica el slide de contenido
    vb.move_slide(prs, len(prs.slides._sldIdLst)-1, 2)
    vb.set_text(prs.slides[1], 3, "Mi título")
    prs.save("salida.pptx")

Por qué clone_slide y no recrear shapes a mano: el template trae layouts ya diseñados
(fondos, decoraciones, mockups). Clonar y reemplazar texto preserva todo el diseño.
"""
import copy
from pptx.oxml.ns import qn


# ---------------------------------------------------------------- introspección
def describe(prs, slide_index=None):
    """Imprime el mapa de slides y shapes. Úsalo SIEMPRE para confirmar índices
    antes de reemplazar texto — son posicionales y cambian si clonas/borras slides."""
    slides = list(prs.slides)
    rng = range(len(slides)) if slide_index is None else [slide_index]
    for i in rng:
        s = slides[i]
        print(f"===== SLIDE idx {i} =====")
        for j, sh in enumerate(s.shapes):
            t = ""
            if sh.has_text_frame and sh.text_frame.text.strip():
                t = " | TXT: " + repr(sh.text_frame.text.strip().replace("\n", " / ")[:55])
            print(f"  [{j}] {str(sh.shape_type):16}{t}")
        print()


# ---------------------------------------------------------------- texto
def set_text(slide, shape_idx, text):
    """Reemplaza el texto de un shape preservando el formato del primer run."""
    tf = slide.shapes[shape_idx].text_frame
    p = tf.paragraphs[0]
    if p.runs:
        p.runs[0].text = text
        for r in p.runs[1:]:
            r.text = ""
    else:
        p.add_run().text = text


def set_paragraphs(slide, shape_idx, texts):
    """Reemplaza una lista de párrafos (bullets, items) preservando formato.
    Si hay más textos que párrafos existentes, ignora el sobrante."""
    tf = slide.shapes[shape_idx].text_frame
    for i, para in enumerate(tf.paragraphs):
        if i < len(texts):
            if para.runs:
                para.runs[0].text = texts[i]
                for r in para.runs[1:]:
                    r.text = ""
            else:
                para.add_run().text = texts[i]
        else:
            for r in para.runs:
                r.text = ""


# ---------------------------------------------------------------- slides
def clone_slide(prs, index):
    """Duplica el slide en `index` (con su fondo, decoraciones e imágenes) y lo
    agrega al FINAL. Devuelve el nuevo slide. Remapea las referencias de imagen
    para que no se pierdan al copiar el XML."""
    source = prs.slides[index]
    new_slide = prs.slides.add_slide(source.slide_layout)
    # limpiar lo que el layout haya inyectado
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)
    # 1) copiar el background del slide (cSld/bg), donde vive el navy/gradiente
    src_cSld = source._element.find(qn("p:cSld"))
    new_cSld = new_slide._element.find(qn("p:cSld"))
    src_bg = src_cSld.find(qn("p:bg"))
    if src_bg is not None:
        new_cSld.insert(0, copy.deepcopy(src_bg))
    # 2) mapear rels de imagen del slide origen al nuevo
    id_map = {}
    for rId, rel in source.part.rels.items():
        if "image" in rel.reltype:
            id_map[rId] = new_slide.part.relate_to(rel._target, rel.reltype)
    # 3) copiar shapes y reapuntar r:embed / r:link a los nuevos rId
    spTree = new_slide.shapes._spTree
    for shp in source.shapes:
        el = copy.deepcopy(shp._element)
        for attr in (qn("r:embed"), qn("r:link")):
            for node in el.iter():
                if attr in node.attrib and node.attrib[attr] in id_map:
                    node.attrib[attr] = id_map[node.attrib[attr]]
        spTree.append(el)
    return new_slide


def move_slide(prs, old_index, new_index):
    """Mueve un slide a otra posición."""
    lst = prs.slides._sldIdLst
    ids = list(lst)
    el = ids[old_index]
    lst.remove(el)
    lst.insert(new_index, el)


def delete_slide(prs, index):
    """Elimina un slide por índice."""
    lst = prs.slides._sldIdLst
    sldId = list(lst)[index]
    rId = sldId.get(qn("r:id"))
    prs.part.drop_rel(rId)
    lst.remove(sldId)


def keep_only(prs, keep_indices):
    """Conserva solo los slides en `keep_indices` (en su orden original) y borra el resto.
    Útil para armar un deck corto a partir del template completo."""
    to_delete = sorted(
        [i for i in range(len(list(prs.slides))) if i not in keep_indices],
        reverse=True,
    )
    for i in to_delete:
        delete_slide(prs, i)
